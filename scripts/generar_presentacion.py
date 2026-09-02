"""Genera presentacion PPT profesional — SIRH Nomina."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colores
PRIMARY = RGBColor(0x1E, 0x6B, 0x8A)
PRIMARY_DARK = RGBColor(0x15, 0x55, 0x70)
ACCENT = RGBColor(0x2A, 0x9D, 0x8F)
DARK = RGBColor(0x1A, 0x23, 0x32)
TEXT_C = RGBColor(0x3D, 0x4F, 0x5F)
MUTED = RGBColor(0x7A, 0x8A, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)
SUCCESS = RGBColor(0x2D, 0x6A, 0x4F)
DANGER = RGBColor(0xC0, 0x39, 0x2B)
WARNING = RGBColor(0xB8, 0x86, 0x0B)
BORDER = RGBColor(0xE2, 0xE7, 0xED)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill_color, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_text(shape, text, size=11, color=TEXT_C, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_text_box(slide, left, top, w, h, text, size=11, color=TEXT_C, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_para(tf, text, size=11, color=TEXT_C, bold=False, space_before=Pt(2), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, size=11, color=TEXT_C, bold_prefix=""):
    p = tf.add_paragraph()
    p.level = 0
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    if bold_prefix:
        run = p.add_run()
        run.text = bold_prefix
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = True
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
    else:
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return p


def add_footer(slide, page_num):
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(8), Inches(0.4),
                 "SIRH Nomina — Cesar Augusto Sotelo Zapata — UCN — Computacion en la Nube",
                 size=8, color=MUTED)
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(0.8), Inches(0.4),
                 str(page_num), size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, left, top, w, rows_data, col_widths):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.4 * rows))
    table = table_shape.table

    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = WHITE if r == 0 else TEXT_C
                paragraph.font.bold = r == 0
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            elif r == len(rows_data) - 1 and rows_data[r][0] == "TOTAL":
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xFA, 0xFB, 0xFC)
    return table


def add_stat_card(slide, left, top, label, value, accent_color):
    w, h = Inches(2.7), Inches(1.5)
    card = add_rounded_rect(slide, left, top, w, h, WHITE)
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)
    # Borde izquierdo
    add_rect(slide, left, top + Inches(0.15), Inches(0.05), h - Inches(0.3), accent_color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.3), w - Inches(0.4), Inches(0.3),
                 label, size=9, color=MUTED, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.65), w - Inches(0.4), Inches(0.6),
                 value, size=28, color=DARK, bold=True)


# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1: PORTADA
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(slide, 0, 0, Inches(13.333), Inches(7.5), DARK)
# Linea accent
add_rect(slide, 0, Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "SIRH Nomina", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.8),
             "Sistema de Gestion de Talento Humano y Liquidacion de Nomina en Pymes",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)

info_lines = [
    ("Curso:", "Computacion en la Nube — UCN"),
    ("Autor:", "Cesar Augusto Sotelo Zapata"),
    ("Metodologia:", "DevOps / CI-CD"),
    ("Infraestructura:", "Linode (Akamai Cloud) — IaaS — USD $7/mes"),
    ("Repositorio:", "github.com/csotelo-dev/ProoyectoCN"),
]
y = 3.8
for label, val in info_lines:
    tf = add_text_box(slide, Inches(4), Inches(y), Inches(5.5), Inches(0.35),
                      "", size=12, color=MUTED, align=PP_ALIGN.LEFT)
    p = tf.paragraphs[0]
    run_l = p.add_run()
    run_l.text = f"{label}  "
    run_l.font.size = Pt(12)
    run_l.font.color.rgb = ACCENT
    run_l.font.bold = True
    run_v = p.add_run()
    run_v.text = val
    run_v.font.size = Pt(12)
    run_v.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    y += 0.35

add_text_box(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.4),
             "Taller ABP — Entrega 1: Ecosistema Interactivo Cloud — Septiembre 2026",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEMA
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "1. Planteamiento del Problema", size=24, color=PRIMARY, bold=True)

# Columna izquierda
tf = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.4),
                  "Situacion actual en Pymes", size=14, color=DARK, bold=True)
tf2 = add_text_box(slide, Inches(0.8), Inches(1.9), Inches(6), Inches(3.6), "", size=11, color=TEXT_C)
bullets_l = [
    ("Errores de liquidacion ", "con impacto economico directo"),
    ("Sin trazabilidad ", "ni control de versiones en datos de nomina"),
    ("Reprocesos administrativos ", "y demoras en atencion de solicitudes"),
    ("Riesgo normativo: ", "incumplimiento Ley 1581/2012 (datos personales)"),
    ("Sin canal de autoconsulta ", "para los empleados"),
    ("Costo fijo on-premise: ", "~USD $120/mes (hardware, energia, licencias)"),
]
for bold_t, normal_t in bullets_l:
    add_bullet(tf2, normal_t, size=11, color=TEXT_C, bold_prefix=f"•  {bold_t}")

# Columna derecha — Stat card
add_stat_card(slide, Inches(7.5), Inches(1.4), "COSTO ON-PREMISE", "$120/mes", DANGER)

# Pregunta problema
box = add_rounded_rect(slide, Inches(7.5), Inches(3.3), Inches(5), Inches(3.0), LIGHT_BG)
box.line.color.rgb = BORDER
box.line.width = Pt(0.5)
add_text_box(slide, Inches(7.8), Inches(3.5), Inches(4.4), Inches(0.3),
             "Pregunta Problema", size=12, color=PRIMARY, bold=True)
add_text_box(slide, Inches(7.8), Inches(4.0), Inches(4.4), Inches(2.0),
             "¿De que manera un sistema de gestion de talento humano y nomina desplegado en un VPS (IaaS) de Linode permite reducir los errores de liquidacion, automatizar los procesos administrativos y disminuir en mas de un 90% el costo fijo mensual de infraestructura?",
             size=10, color=TEXT_C)

add_footer(slide, 2)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 3: SOLUCION
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "2. Solucion Propuesta", size=24, color=PRIMARY, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Sistema web que automatiza la gestion de talento humano y liquidacion de nomina colombiana, desplegado en la nube.",
             size=12, color=TEXT_C)

# 3 stat cards
add_stat_card(slide, Inches(0.8), Inches(2.0), "REDUCCION ERRORES", "-90%", ACCENT)
add_stat_card(slide, Inches(3.8), Inches(2.0), "TIEMPO ATENCION", "-60%", PRIMARY)
add_stat_card(slide, Inches(6.8), Inches(2.0), "AHORRO INFRA", "94.2%", SUCCESS)

# Modulos
y_mod = 3.9
mods = [
    ("Core (Talento Humano)", "CRUD de empleados, departamentos, cargos, contratos. Control de acceso por roles.", PRIMARY),
    ("Payroll (Nomina)", "Motor de liquidacion colombiana: parametros legales, novedades, horas extra, seguridad social, parafiscales.", ACCENT),
    ("API REST (Autoconsulta)", "Endpoints para que empleados consulten desprendibles, contratos y datos personales.", RGBColor(0x34, 0x69, 0x8A)),
]
for name, desc, color in mods:
    card = add_rounded_rect(slide, Inches(0.8), Inches(y_mod), Inches(11.5), Inches(0.8), WHITE)
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)
    add_rect(slide, Inches(0.8), Inches(y_mod), Inches(0.06), Inches(0.8), color)
    add_text_box(slide, Inches(1.1), Inches(y_mod + 0.08), Inches(3), Inches(0.3),
                 name, size=11, color=color, bold=True)
    add_text_box(slide, Inches(1.1), Inches(y_mod + 0.38), Inches(10.5), Inches(0.35),
                 desc, size=10, color=TEXT_C)
    y_mod += 0.95

add_footer(slide, 3)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 4: ARQUITECTURA
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "3. Arquitectura Cloud (IaaS — Linode)", size=24, color=PRIMARY, bold=True)

# VPS Container grande
vps_x, vps_y, vps_w, vps_h = Inches(1.5), Inches(1.5), Inches(10.3), Inches(5.0)
vps_box = add_rounded_rect(slide, vps_x, vps_y, vps_w, vps_h, RGBColor(0xF8, 0xFA, 0xFC))
vps_box.line.color.rgb = PRIMARY
vps_box.line.width = Pt(2)
add_text_box(slide, vps_x + Inches(0.3), vps_y + Inches(0.15), Inches(9), Inches(0.35),
             "Linode Nanode 1GB  —  Ubuntu 24.04 LTS  —  Docker Compose", size=11, color=PRIMARY, bold=True)

# Internet arrow
inet = add_rounded_rect(slide, Inches(0.2), Inches(2.8), Inches(1.2), Inches(0.6), DARK)
set_text(inet, "Internet\nHTTPS", size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.0), Inches(2.95), Inches(0.8), Inches(0.3))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT
arrow.line.fill.background()

# Nginx
nginx_box = add_rounded_rect(slide, Inches(2.0), Inches(2.2), Inches(9.3), Inches(0.9), PRIMARY)
set_text(nginx_box, "Nginx  —  Reverse Proxy + TLS (Let's Encrypt) + Archivos Estaticos", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Django
dj_box = add_rounded_rect(slide, Inches(2.0), Inches(3.4), Inches(9.3), Inches(0.5), RGBColor(0xE8, 0xF4, 0xF8))
dj_box.line.color.rgb = PRIMARY
dj_box.line.width = Pt(1)
set_text(dj_box, "Gunicorn + Django 5  (Python 3.12)", size=11, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# 3 modulos
mod_y = Inches(4.1)
mod_data = [("Core (RRHH)", PRIMARY), ("Payroll (Nomina)", ACCENT), ("API REST", RGBColor(0x34, 0x69, 0x8A))]
mod_x = Inches(2.4)
for name, col in mod_data:
    m = add_rounded_rect(slide, mod_x, mod_y, Inches(2.6), Inches(0.7), col)
    set_text(m, name, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    mod_x += Inches(3.0)

# PostgreSQL
pg_box = add_rounded_rect(slide, Inches(2.0), Inches(5.1), Inches(9.3), Inches(0.8), DARK)
set_text(pg_box, "PostgreSQL 16  (ACID + pgcrypto AES-256)  —  Volumen Persistente Docker", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Seguridad
add_text_box(slide, Inches(2.0), Inches(6.1), Inches(9.3), Inches(0.35),
             "Seguridad: UFW + fail2ban + Linode Cloud Firewall  |  IaC: Terraform  |  CI/CD: GitHub Actions",
             size=9, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(slide, 4)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 5: STACK TECNOLOGICO
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "4. Stack Tecnologico", size=24, color=PRIMARY, bold=True)

data = [
    ["Capa", "Tecnologia", "Justificacion"],
    ["Lenguaje", "Python 3.12 LTS", "Ecosistema maduro para aplicaciones de gestion"],
    ["Framework", "Django 5 + DRF", "ORM, autenticacion, admin y API REST integrados"],
    ["Base de datos", "PostgreSQL 16", "Transacciones ACID; pgcrypto para cifrado AES-256"],
    ["Servidor", "Gunicorn + Nginx", "WSGI produccion + reverse proxy + TLS Let's Encrypt"],
    ["Contenedores", "Docker Compose", "Paridad entre entornos (Twelve-Factor App)"],
    ["IaC", "Terraform (linode)", "Aprovisionamiento declarativo y versionado del VPS"],
    ["CI/CD", "GitHub Actions", "Pipeline: lint, security, test, build, deploy"],
    ["Pruebas", "pytest + Locust", "Tests unitarios, integracion y carga (>= 99% exactitud)"],
    ["Seguridad", "Argon2, bandit", "Hashing OWASP, analisis estatico, auditoria"],
    ["Infraestructura", "Linode Nanode 1GB", "IaaS USD $5/mes + $2 backups = $7/mes total"],
]
add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), data,
          [Inches(2.5), Inches(3), Inches(6.2)])

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
             "Toda la pila es software libre. Alternativas evaluadas: FastAPI (sin admin integrado), MySQL (sin pgcrypto nativo).",
             size=9, color=MUTED)
add_footer(slide, 5)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 6: CI/CD
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "5. DevOps: Pipeline CI/CD e IaC", size=24, color=PRIMARY, bold=True)

# Pipeline visual con flechas
pipeline_steps = [
    ("Lint", "Ruff", PRIMARY),
    ("Security", "Bandit", RGBColor(0x6C, 0x3D, 0x93)),
    ("Test", "pytest + PG", ACCENT),
    ("Build", "Docker", WARNING),
    ("Deploy", "SSH → VPS", SUCCESS),
]
px = Inches(0.8)
for name, tool, col in pipeline_steps:
    box = add_rounded_rect(slide, px, Inches(1.4), Inches(2.1), Inches(1.0), col)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = tool
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    if px > Inches(0.8):
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, px - Inches(0.35), Inches(1.7), Inches(0.35), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        arr.line.fill.background()
    px += Inches(2.45)

# IaC section
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(6), Inches(0.4),
             "Infraestructura como Codigo (Terraform)", size=14, color=DARK, bold=True)

tf = add_text_box(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(3), "", size=11, color=TEXT_C)
add_bullet(tf, "Aprovisiona VPS Nanode 1GB + Linode Cloud Firewall", bold_prefix="•  main.tf:  ")
add_bullet(tf, "Docker, swap 2GB, SSH hardening, fail2ban, UFW", bold_prefix="•  cloud-init:  ")
add_bullet(tf, "Token Linode (sensitive), region, SSH key", bold_prefix="•  variables.tf:  ")

# Deploy flow
add_text_box(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(0.4),
             "Flujo de Despliegue", size=14, color=DARK, bold=True)

tf2 = add_text_box(slide, Inches(7), Inches(3.3), Inches(5.5), Inches(3), "", size=11, color=TEXT_C)
steps = ["git push origin main", "GitHub Actions: lint + security + tests",
         "Build imagen Docker multi-stage", "Deploy via SSH al VPS Linode",
         "docker compose up -d --build", "Migraciones + collectstatic automaticos"]
for i, s in enumerate(steps):
    add_bullet(tf2, s, bold_prefix=f"  {i+1}.  ")

add_footer(slide, 6)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 7: SEGURIDAD
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "6. Seguridad y Cumplimiento Normativo", size=24, color=PRIMARY, bold=True)

# Izquierda: capas
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.35),
             "Defensa en Profundidad", size=13, color=DARK, bold=True)

tf = add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.5), "", size=10, color=TEXT_C)
sec_items = [
    ("Red: ", "Linode Firewall + UFW + fail2ban. SSH ed25519, root deshabilitado."),
    ("Cifrado transito: ", "TLS 1.2+ con Let's Encrypt (renovacion automatica)."),
    ("Cifrado reposo: ", "pgcrypto/AES-256 para campos sensibles de nomina."),
    ("Autenticacion: ", "Contraseñas con Argon2 (OWASP). Min 10 caracteres."),
    ("Autorizacion: ", "Roles con principio de minimo privilegio."),
    ("Auditoria: ", "Middleware registra todas las operaciones de escritura."),
    ("Hardening: ", "CIS Benchmark Ubuntu, parches automaticos."),
    ("Headers HTTP: ", "HSTS, X-Frame-Options, X-Content-Type-Options."),
]
for bold_t, normal_t in sec_items:
    add_bullet(tf, normal_t, size=10, color=TEXT_C, bold_prefix=f"•  {bold_t}")

# Derecha: tabla normativa
data = [
    ["Marco", "Implementacion"],
    ["Ley 1581/2012", "Cifrado datos personales, control acceso, auditoria"],
    ["ISO 27001", "Politica acceso, gestion llaves, monitoreo"],
    ["NIST CSF 2.0", "Identificar, Proteger, Detectar, Responder, Recuperar"],
    ["CIS Benchmark", "Hardening Ubuntu, actualizaciones automaticas"],
    ["OWASP Top 10", "Argon2, CSRF, XSS, SQL injection prevention"],
]
add_table(slide, Inches(7), Inches(1.3), Inches(5.7), data,
          [Inches(2.2), Inches(3.5)])

add_footer(slide, 7)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 8: PRESUPUESTO
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "7. Presupuesto Cloud", size=24, color=PRIMARY, bold=True)

data = [
    ["Servicio", "Proveedor", "Detalle", "USD/mes"],
    ["Nanode 1GB", "Linode (Akamai)", "1 vCPU, 1GB RAM, 25GB SSD, 1TB transfer", "$5.00"],
    ["Backups", "Linode (Akamai)", "Copias automaticas diarias y semanales", "$2.00"],
    ["TOTAL", "", "", "$7.00"],
]
add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), data,
          [Inches(3), Inches(2.5), Inches(4.5), Inches(1.7)])

# Comparativa visual
add_stat_card(slide, Inches(0.8), Inches(3.8), "ON-PREMISE", "$120/mes", DANGER)
add_stat_card(slide, Inches(3.8), Inches(3.8), "LINODE CLOUD", "$7/mes", ACCENT)
add_stat_card(slide, Inches(6.8), Inches(3.8), "AHORRO", "94.2%", SUCCESS)

# TCO
data2 = [
    ["Concepto", "On-Premise", "Linode Cloud", "Ahorro"],
    ["Costo mensual", "$120 USD", "$7 USD", "94.2%"],
    ["Costo anual", "$1,440 USD", "$84 USD", "$1,356 USD"],
]
add_table(slide, Inches(0.8), Inches(5.6), Inches(8), data2,
          [Inches(2.5), Inches(2), Inches(2), Inches(1.5)])

add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
             "Ruta de escalamiento: Linode 2GB ($12/mes) mediante resize sin reinstalacion. Fuente: akamai.com/cloud/pricing",
             size=9, color=MUTED)
add_footer(slide, 8)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 9: RESULTADOS
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "8. Resultados y Pruebas", size=24, color=PRIMARY, bold=True)

data = [
    ["Modulo", "Tests", "Cobertura"],
    ["Core (modelos, vistas, auth)", "7 tests", "CRUD, autenticacion, permisos"],
    ["Payroll (liquidacion nomina)", "9 tests", "Motor de nomina completo"],
    ["API REST (endpoints)", "4 tests", "Permisos, autoconsulta"],
    ["TOTAL", "20 tests", "Exactitud >= 99%"],
]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.7), data,
          [Inches(4), Inches(2.5), Inches(5.2)])

add_text_box(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(0.35),
             "Validaciones del motor de nomina", size=13, color=DARK, bold=True)

tf = add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(2.5), "", size=11, color=TEXT_C)
validations = [
    "Salario proporcional, auxilio de transporte (si salario <= 2 SMMLV)",
    "Horas extra: diurna 25%, nocturna 75%, dominical 100%, nocturna dominical 150%",
    "Aportes empleado: salud 4%, pension 4%  |  Empleador: salud 8.5%, pension 12%",
    "Parafiscales: SENA 2%, ICBF 3%, Caja de Compensacion 4%",
    "Bonificaciones, comisiones, deducciones, prestamos, libranzas",
    "Re-liquidacion sin duplicados. Contratos inactivos excluidos automaticamente.",
]
for v in validations:
    add_bullet(tf, v, size=10, color=TEXT_C, bold_prefix="•  ")

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.3),
             "Repositorio: 72 archivos organizados por capas. README.md con diagramas, instrucciones de despliegue y referencias.",
             size=9, color=MUTED)
add_footer(slide, 9)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 10: CONCLUSIONES
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "9. Conclusiones y Referencias", size=24, color=PRIMARY, bold=True)

# Conclusiones
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.35),
             "Conclusiones", size=14, color=DARK, bold=True)

tf = add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4), "", size=11, color=TEXT_C)
conclusions = [
    ("94.2% ahorro ", "en costos de infraestructura ($120 → $7 USD/mes)"),
    (">= 99% exactitud ", "en liquidacion automatica de nomina"),
    ("60% reduccion ", "en carga administrativa por autoconsulta"),
    ("CI/CD garantiza ", "entregas repetibles y trazables"),
    ("Cumplimiento normativo ", "Ley 1581, ISO 27001, NIST CSF 2.0"),
]
for bold_t, normal_t in conclusions:
    add_bullet(tf, normal_t, size=11, color=TEXT_C, bold_prefix=f"•  {bold_t}")

# Referencias
add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.35),
             "Referencias", size=14, color=DARK, bold=True)

tf2 = add_text_box(slide, Inches(7), Inches(1.7), Inches(5.8), Inches(5), "", size=9, color=MUTED)
refs = [
    "Mell & Grance (2011). NIST SP 800-145.",
    "Kim, Humble, Debois & Willis (2021). The DevOps Handbook.",
    "Morris (2020). Infrastructure as Code. O'Reilly.",
    "Fowler (2006). Continuous Integration.",
    "Wiggins (2017). The Twelve-Factor App.",
    "Pressman & Maxim (2021). Ingenieria del software.",
    "ISO/IEC 25010:2011 — Modelos de calidad.",
    "ISO/IEC 27001:2022 — SGSI.",
    "NIST (2024). Cybersecurity Framework 2.0.",
    "OWASP (2021). OWASP Top 10.",
    "Ley 1581/2012, Decreto 1377/2013 (Colombia).",
    "Akamai (2026). Cloud Pricing — North America.",
]
for r in refs:
    add_bullet(tf2, r, size=9, color=MUTED, bold_prefix="•  ")

add_footer(slide, 10)

# ═════════════════════════════════════════════════════════════════════════
output = "D:/UCN/SEMESTRE 9/COMPUTACIÓN EN LA NUBE/PROYECTOUCN/docs/presentacion.pptx"
prs.save(output)
print(f"PPT generado: {output}")
